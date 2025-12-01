# pptx_processing.py
"""
PPTX Processing Module - Traitement des fichiers PowerPoint

Extraction du contenu des fichiers .pptx:
- Texte des slides (titres, corps, placeholders)
- Notes du présentateur
- Tableaux
- Images avec OCR via LLM Vision

Intégration avec le système OCR LLM existant (llm_ocr.py) pour
extraire le texte des images embarquées dans les présentations.
"""

import os
import re
import io
import logging
from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass, field

logger = logging.getLogger(__name__)

# =============================================================================
#  IMPORT PPTX (avec gestion d'erreur)
# =============================================================================

PPTX_AVAILABLE = False
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    PPTX_AVAILABLE = True
except ImportError:
    logger.warning(
        "[PPTX] python-pptx non installé. "
        "Installez-le avec: pip install python-pptx"
    )

# =============================================================================
#  IMPORT OCR LLM (optionnel)
# =============================================================================

LLM_OCR_AVAILABLE = False
DALLEM_API_KEY = None
DALLEM_API_BASE = None
LLM_MODEL = None
create_http_client = None

try:
    from llm_ocr import (
        image_to_base64,
        ocr_image_with_llm,
        OCRResult,
    )
    from models_utils import (
        DALLEM_API_KEY,
        DALLEM_API_BASE,
        LLM_MODEL,
        create_http_client,
    )
    LLM_OCR_AVAILABLE = True
    logger.info("[PPTX] OCR LLM disponible avec DALLEM")
except ImportError as e:
    logger.debug(f"[PPTX] Module OCR LLM non disponible: {e}")


# =============================================================================
#  DATA CLASSES
# =============================================================================

@dataclass
class SlideContent:
    """Contenu extrait d'une slide."""
    slide_number: int
    title: str = ""
    body_text: str = ""
    notes: str = ""
    table_text: str = ""
    image_texts: List[str] = field(default_factory=list)
    shape_count: int = 0
    image_count: int = 0


@dataclass
class PPTXExtractionResult:
    """Résultat de l'extraction d'un fichier PPTX."""
    filename: str
    total_slides: int
    slides: List[SlideContent]
    full_text: str
    images_processed: int = 0
    images_with_text: int = 0
    ocr_used: bool = False
    processing_time: float = 0.0
    errors: List[str] = field(default_factory=list)


# =============================================================================
#  EXTRACTION DE TEXTE
# =============================================================================

def extract_text_from_shape(shape) -> str:
    """
    Extrait le texte d'une forme PowerPoint.

    Gère les cas: TextFrame, Table, GroupShape.
    """
    texts = []

    try:
        # Texte standard (TextFrame)
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                para_text = ""
                for run in paragraph.runs:
                    if run.text:
                        para_text += run.text
                if para_text.strip():
                    texts.append(para_text.strip())

        # Tableau
        if shape.has_table:
            table_text = extract_text_from_table(shape.table)
            if table_text:
                texts.append(table_text)

        # Groupe de formes (récursif)
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub_shape in shape.shapes:
                sub_text = extract_text_from_shape(sub_shape)
                if sub_text:
                    texts.append(sub_text)

    except Exception as e:
        logger.debug(f"[PPTX] Erreur extraction forme: {e}")

    return "\n".join(texts)


def extract_text_from_table(table) -> str:
    """
    Extrait le texte d'un tableau PowerPoint.

    Format: colonnes séparées par " | ", lignes séparées par newline.
    """
    rows_text = []

    try:
        for row in table.rows:
            cells_text = []
            for cell in row.cells:
                cell_text = ""
                if cell.text_frame:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text:
                                cell_text += run.text
                cells_text.append(cell_text.strip())

            if any(cells_text):
                rows_text.append(" | ".join(cells_text))

    except Exception as e:
        logger.debug(f"[PPTX] Erreur extraction tableau: {e}")

    return "\n".join(rows_text)


def extract_images_from_slide(slide) -> List[Tuple[bytes, str]]:
    """
    Extrait les images d'une slide.

    Returns:
        Liste de tuples (image_bytes, format) ex: (b'...', 'png')
    """
    images = []

    try:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    image_bytes = image.blob
                    # Déterminer le format
                    content_type = image.content_type
                    if 'png' in content_type:
                        fmt = 'png'
                    elif 'jpeg' in content_type or 'jpg' in content_type:
                        fmt = 'jpeg'
                    elif 'gif' in content_type:
                        fmt = 'gif'
                    elif 'bmp' in content_type:
                        fmt = 'bmp'
                    elif 'tiff' in content_type:
                        fmt = 'tiff'
                    else:
                        fmt = 'png'  # Default

                    images.append((image_bytes, fmt))

                except Exception as e:
                    logger.debug(f"[PPTX] Erreur extraction image: {e}")

            # Images dans les groupes
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for sub_shape in shape.shapes:
                    if sub_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            image = sub_shape.image
                            images.append((image.blob, 'png'))
                        except Exception:
                            pass

    except Exception as e:
        logger.debug(f"[PPTX] Erreur parcours shapes: {e}")

    return images


# =============================================================================
#  OCR DES IMAGES
# =============================================================================

def ocr_images_with_llm(
    images: List[Tuple[bytes, str]],
    http_client=None,
    api_key: str = None,
    api_base: str = None,
    model: str = None,
    log=None
) -> List[str]:
    """
    Effectue l'OCR sur une liste d'images via LLM Vision.

    Args:
        images: Liste de (image_bytes, format)
        http_client: Client HTTP (optionnel, créé si absent)
        api_key: Clé API (optionnel, utilise DALLEM par défaut)
        api_base: URL API (optionnel)
        model: Modèle vision (optionnel)
        log: Logger

    Returns:
        Liste de textes extraits (un par image)
    """
    _log = log or logger

    if not LLM_OCR_AVAILABLE:
        print("  ⚠️ [PPTX] OCR LLM non disponible - vérifiez DALLEM_API_KEY")
        _log.warning("[PPTX] OCR LLM non disponible - images ignorées")
        return []

    if not images:
        _log.debug("[PPTX] Aucune image à traiter")
        return []

    print(f"  🔍 [PPTX] OCR de {len(images)} image(s)...")

    # Configuration par défaut (DALLEM)
    if api_key is None:
        api_key = DALLEM_API_KEY
    if api_base is None:
        api_base = DALLEM_API_BASE
    if model is None:
        model = LLM_MODEL
    if http_client is None:
        http_client = create_http_client()

    texts = []

    for i, (image_bytes, fmt) in enumerate(images):
        try:
            print(f"      📷 Image {i+1}/{len(images)} ({fmt}, {len(image_bytes)//1024}KB)...", end=" ", flush=True)
            _log.debug(f"[PPTX] OCR image {i+1}/{len(images)} ({fmt})...")

            text, confidence = ocr_image_with_llm(
                image_bytes=image_bytes,
                http_client=http_client,
                api_key=api_key,
                api_base=api_base,
                model=model,
                log=_log
            )

            if text and confidence > 0.3:
                texts.append(text)
                print(f"✅ {len(text)} chars (conf: {confidence:.0%})")
                _log.debug(f"[PPTX] Image {i+1}: {len(text)} chars, confiance {confidence:.0%}")
            else:
                print(f"⚠️ pas de texte (conf: {confidence:.0%})")
                _log.debug(f"[PPTX] Image {i+1}: pas de texte significatif")

        except Exception as e:
            print(f"❌ Erreur: {e}")
            _log.warning(f"[PPTX] Erreur OCR image {i+1}: {e}")

    print(f"  ✅ [PPTX] OCR terminé: {len(texts)}/{len(images)} images avec texte")
    return texts


# =============================================================================
#  EXTRACTION PRINCIPALE
# =============================================================================

def extract_slide_content(
    slide,
    slide_number: int,
    ocr_images: bool = False,
    http_client=None,
    log=None
) -> SlideContent:
    """
    Extrait le contenu complet d'une slide.

    Args:
        slide: Objet Slide de python-pptx
        slide_number: Numéro de la slide (1-indexed)
        ocr_images: Si True, effectue l'OCR sur les images
        http_client: Client HTTP pour OCR
        log: Logger

    Returns:
        SlideContent avec tout le contenu extrait
    """
    _log = log or logger

    content = SlideContent(slide_number=slide_number)

    # Compteurs
    image_count = 0
    shape_count = 0

    # Textes collectés
    body_texts = []
    table_texts = []

    try:
        for shape in slide.shapes:
            shape_count += 1

            # Titre
            if shape.has_text_frame and shape.is_placeholder:
                try:
                    placeholder_type = shape.placeholder_format.type
                    # Type 1 = TITLE, Type 2 = CENTER_TITLE
                    if placeholder_type in (1, 2, 3):  # Title placeholders
                        if shape.text.strip() and not content.title:
                            content.title = shape.text.strip()
                            continue
                except Exception:
                    pass

            # Images
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_count += 1

            # Tableaux
            if shape.has_table:
                table_text = extract_text_from_table(shape.table)
                if table_text:
                    table_texts.append(table_text)

            # Texte standard
            elif shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    para_text = ""
                    for run in paragraph.runs:
                        if run.text:
                            para_text += run.text
                    if para_text.strip():
                        body_texts.append(para_text.strip())

            # Groupes (récursif)
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for sub_shape in shape.shapes:
                    sub_text = extract_text_from_shape(sub_shape)
                    if sub_text:
                        body_texts.append(sub_text)
                    if sub_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image_count += 1

        # Notes du présentateur
        if slide.has_notes_slide and slide.notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame:
                notes_text = notes_frame.text.strip()
                if notes_text:
                    content.notes = notes_text

        # Assembler
        content.body_text = "\n".join(body_texts)
        content.table_text = "\n".join(table_texts)
        content.shape_count = shape_count
        content.image_count = image_count

        # OCR des images si demandé
        if ocr_images and image_count > 0:
            _log.debug(f"[PPTX] Slide {slide_number}: {image_count} images détectées, extraction...")
            images = extract_images_from_slide(slide)
            _log.debug(f"[PPTX] Slide {slide_number}: {len(images)} images extraites")
            if images:
                image_texts = ocr_images_with_llm(
                    images=images,
                    http_client=http_client,
                    log=_log
                )
                content.image_texts = image_texts
            else:
                _log.warning(f"[PPTX] Slide {slide_number}: {image_count} images comptées mais 0 extraites")
        elif ocr_images:
            _log.debug(f"[PPTX] Slide {slide_number}: aucune image")

    except Exception as e:
        print(f"  ❌ [PPTX] Erreur slide {slide_number}: {e}")
        _log.error(f"[PPTX] Erreur extraction slide {slide_number}: {e}")

    return content


def format_slide_text(content: SlideContent, include_metadata: bool = True) -> str:
    """
    Formate le contenu d'une slide en texte lisible.

    Args:
        content: SlideContent à formater
        include_metadata: Si True, inclut le numéro de slide

    Returns:
        Texte formaté
    """
    parts = []

    if include_metadata:
        parts.append(f"--- Slide {content.slide_number} ---")

    if content.title:
        parts.append(f"# {content.title}")

    if content.body_text:
        parts.append(content.body_text)

    if content.table_text:
        parts.append("\n[Tableau]")
        parts.append(content.table_text)

    if content.notes:
        parts.append("\n[Notes]")
        parts.append(content.notes)

    if content.image_texts:
        parts.append("\n[Texte extrait des images]")
        for i, img_text in enumerate(content.image_texts, 1):
            parts.append(f"Image {i}: {img_text}")

    return "\n".join(parts)


# =============================================================================
#  FONCTIONS PUBLIQUES
# =============================================================================

def extract_text_from_pptx(
    file_path: str,
    ocr_images: bool = False,
    include_notes: bool = True,
    include_tables: bool = True,
    log=None
) -> str:
    """
    Extrait le texte d'un fichier PPTX.

    Args:
        file_path: Chemin vers le fichier .pptx
        ocr_images: Si True, effectue l'OCR sur les images (nécessite LLM)
        include_notes: Si True, inclut les notes du présentateur
        include_tables: Si True, inclut le contenu des tableaux
        log: Logger

    Returns:
        Texte extrait, nettoyé et formaté

    Raises:
        RuntimeError: Si python-pptx n'est pas installé
        FileNotFoundError: Si le fichier n'existe pas
    """
    _log = log or logger

    if not PPTX_AVAILABLE:
        raise RuntimeError(
            "python-pptx n'est pas installé. "
            "Installez-le avec: pip install python-pptx"
        )

    if not os.path.isfile(file_path):
        raise FileNotFoundError(f"Fichier introuvable: {file_path}")

    filename = os.path.basename(file_path)
    _log.info(f"[PPTX] Extraction: {filename}")
    print(f"\n📊 [PPTX] Extraction: {filename}")

    try:
        prs = Presentation(file_path)
    except Exception as e:
        raise RuntimeError(f"Impossible d'ouvrir le fichier PPTX: {e}")

    print(f"  📑 Slides: {len(prs.slides)}")
    print(f"  🔍 OCR images: {'Oui' if ocr_images else 'Non'}")
    print(f"  🤖 LLM OCR disponible: {'Oui' if LLM_OCR_AVAILABLE else 'Non'}")

    # Client HTTP pour OCR (créé une seule fois)
    http_client = None
    if ocr_images:
        if LLM_OCR_AVAILABLE:
            http_client = create_http_client()
            print(f"  ✅ Client HTTP créé pour OCR")
        else:
            print(f"  ⚠️ OCR demandé mais LLM non disponible - vérifiez DALLEM_API_KEY")

    all_texts = []
    total_images = 0
    images_with_text = 0

    for i, slide in enumerate(prs.slides, 1):
        content = extract_slide_content(
            slide=slide,
            slide_number=i,
            ocr_images=ocr_images,
            http_client=http_client,
            log=_log
        )

        # Construire le texte de la slide
        parts = []
        parts.append(f"--- Slide {i} ---")

        if content.title:
            parts.append(f"# {content.title}")

        if content.body_text:
            parts.append(content.body_text)

        if include_tables and content.table_text:
            parts.append(content.table_text)

        if include_notes and content.notes:
            parts.append(f"[Notes: {content.notes}]")

        if content.image_texts:
            for img_text in content.image_texts:
                parts.append(f"[Image: {img_text}]")
                images_with_text += 1

        total_images += content.image_count

        slide_text = "\n".join(parts)
        if slide_text.strip():
            all_texts.append(slide_text)

    # Récapitulatif
    print(f"\n  {'─'*50}")
    print(f"  ✅ [PPTX] Extraction terminée")
    print(f"     📑 Slides: {len(prs.slides)}")
    print(f"     🖼️ Images détectées: {total_images}")
    print(f"     📝 Images avec texte OCR: {images_with_text}")
    print(f"     📄 Caractères extraits: {len(''.join(all_texts)):,}")
    print(f"  {'─'*50}\n")

    _log.info(
        f"[PPTX] Extraction terminée: {len(prs.slides)} slides, "
        f"{total_images} images, {images_with_text} avec texte OCR"
    )

    return "\n\n".join(all_texts)


def process_pptx_with_ocr(
    file_path: str,
    quality_threshold: float = 0.5,
    force_ocr: bool = False,
    log=None
) -> Dict[str, Any]:
    """
    Traite un PPTX avec OCR intelligent sur les images.

    Cette fonction:
    1. Extrait le texte standard (slides, notes, tableaux)
    2. Évalue si les images contiennent probablement du texte
    3. Applique l'OCR LLM sur les images pertinentes

    Args:
        file_path: Chemin vers le fichier .pptx
        quality_threshold: Seuil de qualité pour déclencher l'OCR
        force_ocr: Si True, force l'OCR sur toutes les images
        log: Logger

    Returns:
        Dict avec:
        - text: Texte complet extrait
        - method: "classic" ou "ocr"
        - slides_count: Nombre de slides
        - images_count: Nombre d'images
        - images_ocr: Nombre d'images traitées en OCR
        - ocr_used: True si OCR utilisé
    """
    _log = log or logger

    if not PPTX_AVAILABLE:
        return {
            "text": "",
            "error": "python-pptx non installé",
            "method": "error"
        }

    import time
    start_time = time.time()

    filename = os.path.basename(file_path)
    _log.info(f"[PPTX] Traitement avec OCR: {filename}")

    try:
        prs = Presentation(file_path)
    except Exception as e:
        return {
            "text": "",
            "error": f"Impossible d'ouvrir: {e}",
            "method": "error"
        }

    # Statistiques
    total_slides = len(prs.slides)
    total_images = 0
    images_ocr = 0
    ocr_used = False

    # Déterminer si on doit faire l'OCR
    should_ocr = force_ocr

    if not should_ocr:
        # Compter les images pour décider
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    total_images += 1

        # Si beaucoup d'images, probablement utile de faire l'OCR
        # Heuristique: > 2 images par slide en moyenne
        if total_images > total_slides * 2:
            _log.info(f"[PPTX] {total_images} images détectées - OCR recommandé")
            should_ocr = LLM_OCR_AVAILABLE

    # Extraction
    text = extract_text_from_pptx(
        file_path=file_path,
        ocr_images=should_ocr,
        include_notes=True,
        include_tables=True,
        log=_log
    )

    # Recompter les images OCR si utilisé
    if should_ocr:
        ocr_used = True
        # Compter les "[Image:" dans le texte
        images_ocr = text.count("[Image:")

    processing_time = time.time() - start_time

    result = {
        "text": text,
        "method": "ocr" if ocr_used else "classic",
        "slides_count": total_slides,
        "images_count": total_images,
        "images_ocr": images_ocr,
        "ocr_used": ocr_used,
        "processing_time": processing_time,
    }

    _log.info(
        f"[PPTX] Terminé: {total_slides} slides, "
        f"{total_images} images, OCR={ocr_used}, "
        f"temps={processing_time:.1f}s"
    )

    return result


def extract_pptx_metadata(file_path: str, log=None) -> Dict[str, Any]:
    """
    Extrait les métadonnées d'un fichier PPTX.

    Args:
        file_path: Chemin vers le fichier .pptx
        log: Logger

    Returns:
        Dict avec les métadonnées (titre, auteur, slides, etc.)
    """
    _log = log or logger

    if not PPTX_AVAILABLE:
        return {"error": "python-pptx non installé"}

    try:
        prs = Presentation(file_path)

        # Propriétés du document
        core_props = prs.core_properties

        metadata = {
            "filename": os.path.basename(file_path),
            "slides_count": len(prs.slides),
            "title": core_props.title or "",
            "author": core_props.author or "",
            "subject": core_props.subject or "",
            "keywords": core_props.keywords or "",
            "created": str(core_props.created) if core_props.created else "",
            "modified": str(core_props.modified) if core_props.modified else "",
            "last_modified_by": core_props.last_modified_by or "",
        }

        # Compter les éléments
        total_shapes = 0
        total_images = 0
        total_tables = 0
        slides_with_notes = 0

        for slide in prs.slides:
            for shape in slide.shapes:
                total_shapes += 1
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    total_images += 1
                if shape.has_table:
                    total_tables += 1

            if slide.has_notes_slide and slide.notes_slide:
                if slide.notes_slide.notes_text_frame.text.strip():
                    slides_with_notes += 1

        metadata.update({
            "total_shapes": total_shapes,
            "total_images": total_images,
            "total_tables": total_tables,
            "slides_with_notes": slides_with_notes,
        })

        return metadata

    except Exception as e:
        _log.error(f"[PPTX] Erreur extraction métadonnées: {e}")
        return {"error": str(e)}


# =============================================================================
#  NORMALISATION
# =============================================================================

def normalize_whitespace(text: str) -> str:
    """
    Normalise les espaces dans le texte en conservant les sauts de ligne.
    """
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    lines = text.split("\n")
    norm_lines = []
    for line in lines:
        line = re.sub(r"[ \t]+", " ", line)
        norm_lines.append(line.strip())
    return "\n".join(norm_lines).strip()


# =============================================================================
#  COMPATIBILITÉ AVEC LE PIPELINE
# =============================================================================

def is_pptx_available() -> bool:
    """Vérifie si le module PPTX est disponible."""
    return PPTX_AVAILABLE


def is_pptx_ocr_available() -> bool:
    """Vérifie si l'OCR des images PPTX est disponible."""
    return PPTX_AVAILABLE and LLM_OCR_AVAILABLE
